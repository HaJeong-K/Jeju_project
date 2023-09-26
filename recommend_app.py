# -*- coding:utf-8 -*-

import streamlit as st
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import nltk
from nltk import sent_tokenize, word_tokenize
import re
from konlpy.tag import Kkma
from konlpy.utils import pprint
from wordcloud import WordCloud
from collections import Counter
from PIL import Image, ImageFile
import folium
from folium import plugins
from streamlit_folium import st_folium
import streamlit.components.v1 as components

def run_recommend_app():
    st.header('Home')
    ImageFile.LOAD_TRUNCATED_IMAGES = True
    img_word = wordcloud()
    img_info = info_path()
    # point_dict = vis_app()

    st.markdown('#### 제주 여행지와 관련하여 검색량이 많은 키워드입니다.')
    col1, col2, col3 = st.columns([1,6,1])
    with col2:
        st.image(img_word)
    st.markdown('#### 인기 여행지에 대한 정보입니다.')
    map_vis()
    st.markdown('#### 연도별 제주 여행을 위한 정보 수집 경로 분포입니다.')
    st.image(img_info)
    # st.markdown('#### 2020년-2022년 유형별 인기 방문 구역입니다.')
    # div_dict = {'전체': 0, '제주시': 1, '서귀포시': 2, '남성': 3, '여성': 4,
    #     '20대 미만': 5, '20대-30대': 6, '40대-60대': 7, '70대 이상': 8}
    # div_1 = st.radio('구분', ['구역별', '성별', '연령별'], horizontal = True)
    # if div_1 == '구역별':
    #     div_2 = st.radio('구역을 선택해주세요.', ['전체', '제주시', '서귀포시'], horizontal = True)
    # elif div_1 == '성별':
    #     div_2 = st.radio('성별을 선택해주세요.', ['남성', '여성'], horizontal = True)
    # elif div_1 == '연령별':
    #     div_2 = st.radio('연령대를 선택해주세요.', ['20대 미만', '20대-30대', '40대-60대', '70대 이상'], horizontal = True)
    # else:
    #     pass
    # idx = div_dict[div_2]
    # col4, col5, col6 = st.columns([1,6,1])
    # with col5:
    #     st.image(point_dict[idx])

        # images = []
        # index = st.slider('figure', 0, 10)
        # st.image(images[])
        
        # st.image('./data/월간보고서/image/11_202212_1.png')
    
def wordcloud():
    nltk.download('punkt')
    trip_advisor = '한림공원, 비자림, 제주 돌문화 공원, 제주 절물 자연 휴양림, 제주 4.3 평화공원, 한라 수목원, 우도등대공원, 사라봉공원, 함덕 서우봉 해변, 협재 해수욕장, 월정리 해변, 세화 해수욕장, 이호테우 해변, 금능해수욕장, 김녕 성세기 해변, 곽지 해수욕장, 에코랜드 테마파크, 제주 유리의 성, 테지움 제주, 더 마 파크, 제주 공룡 랜드, 제주미니미니랜드, 파크 써던랜드, 스누피 가든, 한림공원, 제주김녕미로공원, 메이즈랜드, 다희연, 김경숙해바라기농장, 제주불빛정원, 카페 글렌코, 토끼섬, 수월봉, 제주김녕미로공원, 신비의 도로, 용두암, 신창 풍차해안도로, 항파두리 항몽, 제주불빛정원, 수월봉, 한담해안산책로, 신창 풍차해안도로, 검멀레, 서우봉, 관음사탐방로, 왕따나무, 민오름 - 오라동, 녹산로유채꽃길, 우도, 만장굴, 함덕 서우봉 해변, 동문 재래시장, 협재 해수욕장, 한림공원, 에코랜드 테마파크, 비자림, 산굼부리 분화구, 한라산 성판악. 해녀 박물관, 성읍 민속마을 '
    hotels = '한라산, 성산일출봉, 섭지코지, 만장굴, 우도, 추자도, 국립제주박물관, 성읍민속마을, 휴애리 자연생활공원, 제주 올레'
    jeju_iljung = '섭지코지, 용두암&용연, 오설록티뮤지엄, 사계해안, 하귀↔애월 해안도로, 쇠소깍 , 한담해안산책로 , 월정리해변 , 제주시 동문재래시장 , 신비의도로(경유) , 협재해수욕장(산책) , 외돌개 , 새연교 , 1100고지습지 , 서귀포매일올레시장'
    trip = '제주특별자치도 민속 자연사 박물관, 브릭캠퍼스 제주, 제주해녀박물관, 국립제주박물과느, 삼성혈, 수목원테마파크, 제주민속박물관, 테지움 제주, 아라리오뮤지엄, 메이즈랜드, 제주현대미술관'
    trip2 = '우도, 한라산 국립공원, 용두암, 에코랜드, 산굼부리, 월정리 갈비밥, 만장굴, 한림공원, 용연구름다리, 김녕 미로공원, 제주 돌문화공원, 사려니숲길, 백록담, 거문오름'
    expedia = '섭지코지, 방주 교회, 성이시돌 목장, 김경숙 해바라기 농장, 이중섭 문화거리, 돈내코 원앙폭포, 오설록 티 뮤지엄, 신천 목장'
    visit_jeju = '우도, 성산일출봉, 사려니숲길, 카멜리아힐, 협재해수욕장, 월정리해변, 함덕해수욕장, 산굼부리, 새별오름, 섭지코지, 오설록티뮤지엄, 한라산국립공원, 제주4.3평화공원, 에코랜드테마파크, 용눈이오름, 비자림, 천지연폭포, 세화해변, 우도산호해변 홍조단괴 서빈백사, 광치기해변, 김녕해수욕장, 아쿠아플라넷 제주, 송악산, 쇠소깍, 금능 해수욕장, 마라도, 애월해안도로, 곽지해수욕장, 용머리해안, 한라산영실코스'
    news = '동문재래시장, 서귀포매일올레시장, 성산일출봉, 함덕해수욕장, 협재해수욕장, 이호테우해수욕장, 아르떼뮤지엄 제주, 섭지코지, 곽지해수욕장, 비자림, 금능해수욕장, 김녕해수욕장, 천지연폭포, 금오름, 제주민속오일시장, 수목원길야시장, 사려니숲길, 제주민속촌'
    google = '협재해변, 함덕해수욕장, 만장굴, 월정리 해수욕장, 한림공원, 산굼부리, 에코랜드, 오설록 티 뮤지엄, 우도, 섭지코지, 비자림, 한라수목원'
    top_10 = '외돌개, 천제연폭포, 성산일출봉, 용머리해안, 주상절리대, 정방폭포, 에코랜드, 섭지코지, 용두암, 오설록 티 뮤지엄'

    trip_adv_tokens = Seq1(trip_advisor).process_string().tokens
    hotels_tokens = Seq1(hotels).process_string().tokens
    iljung_tokens = Seq1(jeju_iljung).process_string().tokens
    trip_tokens = Seq1(trip).process_string().tokens
    trip2_tokens = Seq1(trip2).process_string().tokens
    expedia_tokens =Seq1(expedia).process_string().tokens
    visit_jeju_tokens = Seq1(visit_jeju).process_string().tokens
    news_tokens = Seq1(news).process_string().tokens
    google_tokens = Seq1(google).process_string().tokens
    top_10_tokens = Seq1(top_10).process_string().tokens

    words = trip_adv_tokens+hotels_tokens+iljung_tokens+trip_tokens+trip2_tokens+expedia_tokens+visit_jeju_tokens+news_tokens+google_tokens+top_10_tokens
    words = [word for word in words if word not in ['(', ')', '&']]
    counts = Counter(words)
    tags = counts.most_common(60)

    im = Image.open('./data/jeju.png') # 이미지 파일 읽어오기
    FONT_PATH = './customFonts/NanumGothic-Regular.ttf'
    mask_arr = np.array(im) # 픽셀 값 배열 형태 변환
    background_color = (190, 247, 255)
    wc = WordCloud(font_path = FONT_PATH, background_color = background_color,
                max_font_size = 60, width = 1010, height = 800, mask = mask_arr,
                prefer_horizontal = True, colormap = 'gist_earth_r')
    cloud = wc.generate_from_frequencies(dict(tags))

    fig, ax = plt.subplots(figsize = (15, 15))
    ax.set_axis_off()
    ax.imshow(cloud, interpolation = 'bilinear')
    fig.savefig('./data/streamlit/wordcloud.png', bbox_inches = 'tight')
    image = Image.open('./data/streamlit/wordcloud.png')
    
    return image

class Seq1:
    def __init__(self, input_string):
        self.input_string = input_string
        self.tokens = None  # 토큰 리스트를 저장할 변수 초기화

    def remove_spaces(self):
        self.input_string = self.input_string.replace(' ', '')

    def remove_commas(self):
        self.input_string = self.input_string.replace(',', ' ')

    def remove_some(self):
        self.input_string = self.input_string.replace('국립공원', ' ')

    def tokenize_words(self):
        self.tokens = nltk.word_tokenize(self.input_string)  # 토큰 리스트 저장
        self.input_string = ' '.join(self.tokens)

    def process_string(self):
        self.remove_spaces()
        self.remove_commas()
        self.remove_some()
        self.tokenize_words()

        return self

def map_vis():
    m = folium.Map(
        location = [33.361936, 126.529165], # 코드 실행시 보이는 제주도의 좌표
        zoom_start = 11                       # 코드 실행시 보이는 지도 크기 조절
    )

    # 공원 & 테마파크
    create_marker(m, [33.389644, 126.239229], "한림 공원", "http://www.hallimpark.com/", "fa-tree", "green")
    create_marker(m, [33.537020, 126.772977], "김녕미로공원", "http://www.jejumaze.com/", "fa-route", "green")
    create_marker(m, [33.422436, 126.409554], "불빛정원", "http://www.rosestar.kr/", "fa-wand-magic-sparkles", "beige")
    create_marker(m, [33.455789, 126.668477], "제주 에코랜드", "https://www.ecolandjeju.co.kr/", "fa-cannabis", "green")
    create_marker(m, [33.448709, 126.659280], "제주 돌문화 공원", "http://www.jeju.go.kr/jejustonepark/index.htm", "fa-cannabis", "green")
    create_marker(m, [33.494462, 126.964894], "우도 등대 공원", "https://www.tripadvisor.co.kr/Attraction_Review-g297885-d4094868-Reviews-Udo_Island_Lighthouse_Park-Jeju_Jeju_Island.html", "fa-tower-cell", "#FF0000")
    create_marker(m, [33.517937, 126.544849], "사라봉 공원", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500277", "fa-cannabis", "green")
    create_marker(m, [33.444340, 126.778274], "스누피 가든", "https://www.instagram.com/snoopygardenkorea/", "fa-cannabis", "green")
    create_marker(m, [33.488867, 126.800951], "메이즈랜드", "https://mazeland.co.kr/kor/", "fa-route", "green")
    create_marker(m, [33.451773, 126.489103], "러브랜드", "http://www.jejuloveland.com/", "fa-cannabis", "pink")
    create_marker(m, [33.451785, 126.618972], "4·3 평화공원", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500535", "fa-leaf", "green")
    create_marker(m, [33.439658, 126.629467], "절물자연휴양림", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500570", "fa-tree", "green")
    create_marker(m, [33.470061, 126.493234], "한라 수목원", "https://www.jeju.go.kr/sumokwon/index.htm", "fa-tree", "green")
    create_marker(m, [33.308828, 126.633740], "휴애리 자연생활공원", "http://m.hueree.com/index.php", "fa-leaf", "green")
    create_marker(m, [33.357727, 126.463078], "1100고지 습지", "https://www.visitjeju.net/kr/detail/view?contentsid=CNTS_000000000021240", "fa-leaf", "green")
    create_marker(m, [33.440995, 126.434542], "제주 공룡랜드", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500536", "fa-landmark", "beige")
    create_marker(m, [33.458436, 126.485446], "브릭캠퍼스 제주", "http://www.brickcampus.com/page/exhibition?gbn=1", "fa-puzzle-piece", "green")
    create_marker(m, [33.314685, 126.273709], "유리의 성", "http://www.jejuglasscastle.com/", "fa-gopuram", "lightblue")
    create_marker(m, [33.412365, 126.393523], "테지움 제주", "http://www.teseum.net/default/sub2/sub22.php", "fa-landmark", "purple")
    create_marker(m, [33.433733, 126.673039], "제주 센트럴파크", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500552", "fa-cannabis", "green")
    create_marker(m, [33.445740, 126.547590], "제주 난타 전용관", "https://kr.trip.com/travel-guide/attraction/jeju/jeju-nanta-10759495/", "fa-eye", "black")
    create_marker(m, [33.355968, 126.242346], "더마파크", "http://www.mapark.co.kr/", "fa-horse", "beige")
    create_marker(m, [33.470591, 126.488222], "수목원 테마파크", "http://www.sumokwonpark.com/", "fa-landmark", "green")
    create_marker(m, [33.452197, 126.407715], "항파두리 항몽유적지", "http://www.jeju.go.kr/hangpadori/index.htm", "fa-landmark", "black")
    create_marker(m, [33.484704, 126.806517], "비자림", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500270", "fa-tree", "green")

    # 해수욕장 & 폭포
    create_marker(m, [33.543577, 126.669668], "함덕 해수욕장", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500693", "fa-umbrella-beach", "blue")
    create_marker(m, [33.394429, 126.239924], "협재 해수욕장", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500697", "fa-umbrella-beach", "blue")
    create_marker(m, [33.343516, 126.174124], "신창 풍차해안", "https://www.visitjeju.net/kr/detail/view?contentsid=CNTS_200000000007676", "fa-umbrella-beach", "blue")
    create_marker(m, [33.556307, 126.796126], "월정리 해수욕장", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500496", "fa-umbrella-beach", "blue")
    create_marker(m, [33.525572, 126.860152], "세화해변", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500361", "fa-umbrella-beach", "blue")
    create_marker(m, [33.498088, 126.453159], "이호테우 해변", "https://www.tripadvisor.co.kr/Attraction_Review-g297885-d10586744-Reviews-Iho_Tewoo_Beach-Jeju_Jeju_Island.html", "fa-umbrella-beach", "blue")
    create_marker(m, [33.390178, 126.235278], "금능 해수욕장", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500079", "fa-umbrella-beach", "blue")
    create_marker(m, [33.451062, 126.305740], "곽지 해수욕장", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500056", "fa-umbrella-beach", "blue")
    create_marker(m, [33.459228, 126.310401], "한담 해안산책로", "https://www.visitjeju.net/kr/detail/view?contentsid=CNTS_000000000020159", "fa-shoe-prints", "blue")
    create_marker(m, [33.497391, 126.967168], "검멀레 해변", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500046", "fa-umbrella-beach", "blue")
    create_marker(m, [33.558250, 126.759841], "김녕 해수욕장", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500083", "fa-umbrella-beach", "blue")
    create_marker(m, [33.231259, 126.310443], "사계해변", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500274", "fa-umbrella-beach", "blue")
    create_marker(m, [33.299809, 126.582286], "돈내코 원앙폭포", "https://www.visitjeju.net/kr/detail/view?contentsid=CNTS_000000000019594", "fa-water", "blue")
    create_marker(m, [33.246939, 126.554400], "천지연 폭포", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500618&menuId=DOM_000001718000000000#p2", "fa-water", "blue")
    create_marker(m, [33.502758, 126.943260], "서빈백사 해변(산호해변)", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500474&menuId=DOM_000001718000000000#p2", "fa-umbrella-beach", "blue")
    create_marker(m, [33.452447, 126.924673], "광치기 해변", "https://www.visitjeju.net/kr/detail/view?contentsid=CNTS_000000000018413&menuId=DOM_000001718000000000#p2", "fa-umbrella-beach", "blue")
    create_marker(m, [33.231874, 126.314641], "용머리 해안", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500471&menuId=DOM_000001718000000000#p2", "fa-eye", "black")
    create_marker(m, [33.237762, 126.425093], "주상절리대", "https://www.visitjeju.net/kr/detail/view?contentsid=CNTS_000000000020476", "fa-eye", "black")
    create_marker(m, [33.244855, 126.571823], "정방폭포", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500531", "fa-water", "blue")
    create_marker(m, [33.516335, 126.512058], "용두암", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500469", "fa-umbrella-beach", "blue")
    create_marker(m, [33.514892, 126.514372], "용연", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500472", "fa-bridge", "black")
    create_marker(m, [33.252550, 126.623523], "쇠소깍", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500384", "fa-water", "blue")

    # 예술 및 관람
    create_marker(m, [33.432846, 126.927833], "아쿠아플라넷 제주", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500565&menuId=DOM_000001718000000000#p2", "fa-fish-fins", "blue")
    create_marker(m, [33.506560, 126.531624], "제주 민속자연사박물관", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500750", "fa-landmark", "black")
    create_marker(m, [33.523547, 126.863489], "해녀박물관", "http://www.jeju.go.kr/haenyeo/index.htm", "fa-landmark", "black")
    create_marker(m, [33.513600, 126.548901], "국립제주박물관", "https://jeju.museum.go.kr/html/kr/index.html", "fa-landmark", "black")
    create_marker(m, [33.305902, 126.289566], "오설록 티 뮤지엄", "https://www.osulloc.com/", "fa-mug-hot", "purple")
    create_marker(m, [33.386361, 126.799689], "성읍민속마을", "http://www.jeju.go.kr/seongeup/index.htm", "fa-house-chimney", "black")
    create_marker(m, [33.517431, 126.523436], "아라리오 뮤지엄", "http://www.arariomuseum.org/exhibition/#/cinema.php", "fa-landmark", "black")
    create_marker(m, [33.339574, 126.265633], "제주현대미술관", "https://www.jeju.go.kr/jejumuseum/index.htm", "fa-landmark", "black")
    create_marker(m, [33.396494, 126.344940], "아르떼 뮤지엄 제주", "https://artemuseum.com/JEJU", "fa-landmark", "black")
    create_marker(m, [33.245623, 126.564222], "이중섭 문화거리", "https://www.visitjeju.net/kr/detail/view?contentsid=CNTS_000000000020199", "fa-road", "black")
    create_marker(m, [33.237488, 126.559781], "새연교", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500312", "fa-bridge", "black")
    create_marker(m, [33.504697, 126.529340], "삼성혈", "http://samsunghyeol.or.kr/index.php", "fa-landmark", "black")
    create_marker(m, [33.344697, 126.326672], "성이시돌목장", "https://www.visitjeju.net/kr/detail/view?contentsid=CNTS_200000000008053", "fa-horse", "beige")
    create_marker(m, [33.305111, 126.387428], "방주교회", "https://www.visitjeju.net/kr/detail/view?contentsid=CNTS_000000000018384", "fa-church", "black")

    # 쇼핑 & 시장
    create_marker(m, [33.492912, 126.493513], "BAOJIAN STREET(누웨마루 거리)", "https://kr.trip.com/travel-guide/attraction/jeju/nuwemaru-street-94751/", "fa-cart-shopping", "blue")
    create_marker(m, [33.513584, 126.524534], "칠성로 쇼핑거리", "https://www.jeju.go.kr/sumokwon/index.htm", "fa-cart-shopping", "blue")
    create_marker(m, [33.515431, 126.526859], "흑돼지 거리", "https://www.tripadvisor.co.kr/Attraction_Review-g297885-d8029124-Reviews-Black_Pork_Street-Jeju_Jeju_Island.html", "fa-piggy-bank", "pink")
    create_marker(m, [33.512049, 126.528280], "동문 재래시장", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500745", "fa-cart-shopping", "blue")
    create_marker(m, [33.248708, 126.564102], "서귀포 매일올레시장", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500731", "fa-cart-shopping", "blue")
    create_marker(m, [33.496964, 126.475372], "제주민속오일 시장", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500748", "fa-cart-shopping", "blue")
    create_marker(m, [33.470218, 126.488197], "수목원길 야시장", "http://www.sumokwonpark.com/theme/case2/sub/040201.php", "fa-cart-shopping", "blue")
    create_marker(m, [33.433608, 126.733369], "카페 글렌코", "https://www.instagram.com/cafe_glencoe/", "fa-mug-hot", "purple")

    # 오름 & 자연 경관 & 산
    create_marker(m, [33.478948, 126.627650], "김경숙 해바라기농장", "https://www.visitjeju.net/kr/detail/view?contentsid=CNTS_000000000018385", "fa-leaf", "green")
    create_marker(m, [33.420318, 126.549048], "한라산 관음사 탐방로", "https://visithalla.jeju.go.kr/contents/contents.do?id=62", "fa-shoe-prints", "green")
    create_marker(m, [33.360896, 126.535767], "한라산 국립공원", "http://www.jeju.go.kr/hallasan/index.htm", "fa-tree", "green")
    create_marker(m, [33.478518, 126.502606], "민오름", "https://www.visitjeju.net/kr/detail/view?contentsid=CNTS_000000000021534", "fa-shoe-prints", "green")
    create_marker(m, [33.398812, 126.720037], "녹산로 유채꽃길", "https://www.visitjeju.net/kr/detail/view?contentsid=CNTS_000000000000975", "fa-leaf", "green")
    create_marker(m, [33.458306, 126.942586], "성산 일출봉", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500349", "fa-mountain-sun", "green")
    create_marker(m, [33.454687, 126.717914], "거문오름", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500034", "fa-mountain-sun", "green")
    create_marker(m, [33.393697, 126.682347], "사려니숲", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500281", "fa-tree", "green")
    create_marker(m, [33.356635, 126.868217], "신천목장", "https://www.tripadvisor.co.kr/Attraction_Review-g297892-d9710232-Reviews-Sinpung_Sincheon_Bada_Mokjang-Seogwipo_Jeju_Island.html", "fa-eye", "black")
    create_marker(m, [33.289823, 126.368298], "카멜리아 힐", "https://www.visitjeju.net/kr/detail/view?contentsid=CNTS_000000000001195", "fa-leaf", "green")
    create_marker(m, [33.460235, 126.831477], "용눈이오름", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500466", "fa-mountain-sun", "green")
    create_marker(m, [33.356209, 126.307619], "금오름", "https://www.visitjeju.net/kr/detail/view?contentsid=CNTS_200000000012906", "fa-mountain-sun", "green")
    create_marker(m, [33.351556, 126.350339], "새별오름 나홀로나무", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500309", "fa-tree", "green")
    create_marker(m, [33.197902, 126.291922], "송악산", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500378&menuId=DOM_000001718000000000#p2", "fa-mountain-sun", "green")
    create_marker(m, [33.451209, 126.487808], "신비의도로", "https://www.tripadvisor.co.kr/Attraction_Review-g297885-d1747359-Reviews-Mysterious_Road_Dokkaebi_Road-Jeju_Jeju_Island.html", "fa-road", "black")
    create_marker(m, [33.528291, 126.770309], "만장굴", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500182", "fa-dungeon", "black")
    create_marker(m, [33.545737, 126.677807], "서우봉", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500335", "fa-mountain-sun", "green")
    create_marker(m, [33.295375, 126.164100], "수월봉", "https://www.visitjeju.net/kr/detail/view?contentsid=CNTS_000000000020148", "fa-shoe-prints", "blue")
    create_marker(m, [33.431609, 126.690000], "산굼부리", "http://www.sangumburi.net/", "fa-leaf", "green")
    create_marker(m, [33.487446, 126.706122], "제주 다희연", "https://www.tripadvisor.co.kr/Attraction_Review-g297885-d2337358-Reviews-DaHeeYeon-Jeju_Jeju_Island.html", "fa-leaf", "green")
    create_marker(m, [33.424227, 126.931111], "섭지코지", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500343", "fa-shoe-prints", "blue")
    create_marker(m, [33.240042, 126.545719], "외돌개", "https://www.visitjeju.net/kr/detail/view?contentsid=CNTS_000000000018409", "fa-eye", "black")

    # 섬
    create_marker(m, [33.524438, 126.902811], "토끼섬", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500652", "fa-circle-info", "#FF0000")
    create_marker(m, [33.947149, 126.319924], "추자도", "https://www.visitjeju.net/kr/detail/view?contentsid=CNTS_000000000018441", "fa-circle-info", "#FF0000")
    create_marker(m, [33.504511, 126.954134], "우도", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500477", "fa-circle-info", "#FF0000")
    create_marker(m, [33.120649, 126.267375], "마라도", "https://www.visitjeju.net/kr/detail/view?contentsid=CONT_000000000500171&menuId=DOM_000001718000000000#p2", "fa-circle-info", "#FF0000")

    return folium_static(m)

def create_marker(m, location, name, link = None, icon = None, icon_color = 'black'):
    popup_content = f'<div style="width:100px"><strong>{name}</strong><br>' # 관광지 이름 CSS 설정
    if link:
        popup_content += f'<a href="{link}" target="_blank">상세 페이지 확인</a>' # A태그 활용 관광지 클릭시 상세페이지 연결 창이 뜨도록 설정
    popup_content += ''

    folium.Marker(
        location = location,   # 관광지 좌표를 설정해서 지도에 표시
        popup = popup_content, # 관광지 이름 설정
        tooltip = name,        # 관광지에 마우스를 올려놓을 시 보이는 이름
        icon = folium.Icon(color = icon_color, prefix = 'fa-solid' if not icon else 'fa', icon = icon if icon else '') # 관광지에 대한 간단한 아이콘 표시
    ).add_to(m)

def folium_static(fig, width=700, height=500):
    if isinstance(fig, folium.Map):
        fig = folium.Figure().add_child(fig)
        return components.html(
            fig.render(), height=(fig.height or height) + 10, width=width
            )

    elif isinstance(fig, plugins.DualMap):
        return components.html(
            fig._repr_html_(), height=height + 10, width=width
        )

def data_preprocessing():
    # 2020년 1월부터 2022년 12월까지 YYYYMM 형태로 리스트에 저장
    YYYYMM = []
    for i in [2020, 2021, 2022]:
        for j in range(1, 13):
            YYYYMM.append(f'{i}{j:02d}')

    slides = [10, 20, 32, 51, 52, 55, 56, 57, 58] # 0-indexing 되어있는 pptx파일 슬라이드 중 가공하고자 하는 인덱스

    DATA_PATH = './data/월간보고서/'

    df_dict = {}

    for t in YYYYMM:
        parsed = Presentation(DATA_PATH + t + '.pptx')  # pptx파일 불러오기
        for page in slides:
            slide = parsed.slides[page] # pptx파일에서 가공하고자 하는 슬라이드 순회하며 불러오기
            idx_1 = 0 # page별로 table이 여러 개 있는 경우 구분을 위해 별도 idx 설정
            idx_2 = 0
            for shape in slide.shapes:
                table_data = []
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:  # 슬라이드 내 객체의 type이 table인 경우 추출
                    idx_1 += 1
                    row_count = len(shape.table.rows)
                    col_count = len(shape.table.columns)
                    for r in range(0, row_count):
                        row = []
                        for c in range(0, col_count):
                            cell = shape.table.cell(r, c)
                            row.append(cell.text) # table의 (r, c)위치에 존재하는 text를 row 리스트에 저장
                        table_data.append(row)  # r과 c를 순회하며 row 리스트를 table_data 리스트에 저장하여 2d-array 생성
                    df_temp = pd.DataFrame(table_data)
                    df_dict[f'{t}_{page+1:02d}_{idx_1}'] = df_temp  # df_dict에 '시기_슬라이드번호_idx'를 key로 가지도록 dataframe 저장
                # elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                #     idx_2 += 1
                #     image_blob = shape.image.blob
                #     ext = shape.image.ext
                #     with open(f'{DATA_PATH}image/{page+1:02d}_{t}_{idx_2}.{ext}', 'wb') as file:
                #         file.write(image_blob)
                else:
                    pass

    slides_1 = [10, 20, 32, 51, 52]

    pre_df_dict = {}

    for t in YYYYMM:
        for page in slides:
            if page in slides_1:
                df_temp = df_dict[f'{t}_{page+1:02d}_1']
                if page == 10:
                    df_temp = pre_total(df_temp, t)
                    pre_df_dict[f'{t}_{page+1:02d}'] = df_temp
                elif page == 20 or 32:
                    df_temp = pre_partial(df_temp, t)
                    pre_df_dict[f'{t}_{page+1:02d}'] = df_temp
                elif page == 51:
                    df_temp = pre_male(df_temp, t)
                    pre_df_dict[f'{t}_{page+1:02d}'] = df_temp
                else:
                    df_temp = pre_female(df_temp, t)
                    pre_df_dict[f'{t}_{page+1:02d}'] = df_temp
            else:
                df_temp = df_dict[f'{t}_{page+1:02d}_2']
                df_temp = pre_age(df_temp, t)
                pre_df_dict[f'{t}_{page+1:02d}'] = df_temp

    total_df = {}

    visit_ratio_total = []
    visit_ratio_jejusi = []
    visit_ratio_seogwiposi = []
    visit_count_male = []
    visit_count_female = []
    visit_count_below20 = []
    visit_count_2030 = []
    visit_count_4060 = []
    visit_count_above70 = []

    # 전 기간 데이터를 하나의 dataframe으로 병합
    for t in YYYYMM:
        for page in slides:
            df_temp = pre_df_dict[f'{t}_{page+1:02d}']
            if page == 10:
                visit_ratio_total.append(df_temp)
            elif page == 20:
                visit_ratio_jejusi.append(df_temp)
            elif page == 32:
                visit_ratio_seogwiposi.append(df_temp)
            elif page == 51:
                visit_count_male.append(df_temp)
            elif page == 52:
                visit_count_female.append(df_temp)
            elif page == 55:
                visit_count_below20.append(df_temp)
            elif page == 56:
                visit_count_2030.append(df_temp)
            elif page == 57:
                visit_count_4060.append(df_temp)
            else:
                visit_count_above70.append(df_temp)

    ratio_total = pd.concat(visit_ratio_total).reset_index(drop = True)
    ratio_jejusi = pd.concat(visit_ratio_jejusi).reset_index(drop = True)
    ratio_seogwiposi = pd.concat(visit_ratio_seogwiposi).reset_index(drop = True)
    count_male = pd.concat(visit_count_male).reset_index(drop = True)
    count_female = pd.concat(visit_count_female).reset_index(drop = True)
    count_below20 = pd.concat(visit_count_below20).reset_index(drop = True)
    count_2030 = pd.concat(visit_count_2030).reset_index(drop = True)
    count_4060 = pd.concat(visit_count_4060).reset_index(drop = True)
    count_above70 = pd.concat(visit_count_above70).reset_index(drop = True)

    return ratio_total, ratio_jejusi, ratio_seogwiposi, count_male, count_female, count_below20, count_2030, count_4060, count_above70

def pre_total(df, t):
    df = df.loc[2:16, [3, 7]]
    df.reset_index(drop = True, inplace = True)
    df.reset_index(inplace = True)
    df.rename(columns = {'index': '순위', 3: '블록명', 7: '비율'}, inplace = True)
    df['순위'] = df['순위'] + 1
    df['시기'] = t
    return df

def pre_partial(df, t):
    df = df.loc[2:16, [2, 6]]
    df.reset_index(drop = True, inplace = True)
    df.reset_index(inplace = True)
    df.rename(columns = {'index': '순위', 2: '블록명', 6: '비율'}, inplace = True)
    df['순위'] = df['순위'] + 1
    df['시기'] = t
    return df

def pre_male(df, t):
    df = df.loc[2:, [0, 2, 3]]
    df.reset_index(drop = True, inplace = True)
    df.rename(columns = {0: '순위', 2: '블록명', 3: '남성_방문객_수'}, inplace = True)
    df['순위'] = df['순위'].astype('int64')
    df['시기'] = t
    return df

def pre_female(df, t):
    df = df.loc[2:, [0, 2, 5]]
    df.reset_index(drop = True, inplace = True)
    df.rename(columns = {0: '순위', 2: '블록명', 5: '여성_방문객_수'}, inplace = True)
    df['순위'] = df['순위'].astype('int64')
    df['시기'] = t
    return df

def pre_age(df, t):
    div = df.loc[0, 3]
    df = df.loc[2:, [0, 2, 3]]
    df.reset_index(drop = True, inplace = True)
    df.rename(columns = {0: '순위', 2: '블록명', 3: f'{div} 방문객 수'}, inplace = True)
    df['순위'] = df['순위'].astype('int64')
    df['시기'] = t
    return df

# def vis_app():
#     point_dict = {}
#     for idx in range(9):
#         plt.rcParams['font.family'] = 'NanumGothic'
#         fig, ax = plt.subplots(figsize = (15, 5))
#         temp_df = data_preprocessing()[idx]
#         counts = temp_df['블록명'].value_counts()
#         pop_sec = counts[counts >= 25].index.to_list()
#         temp_df = temp_df.loc[temp_df['블록명'].isin(pop_sec), :]
#         sns.pointplot(data = temp_df, x = '시기', y = '순위', hue = '블록명', ax = ax)
#         ax.set_xticklabels(ax.get_xticklabels(), rotation = -45)
#         ax.legend(loc = 'center left', bbox_to_anchor = (1, 0.5))
#         ax.invert_yaxis()
#         ax.grid(False)
#         fig.savefig(f'./data/streamlit/pointplot_{idx}.png', bbox_inches = 'tight')
#         image = Image.open(f'./data/streamlit/pointplot_{idx}.png')
#         point_dict[idx] = image

#     return point_dict

def info_path():
    info_dict = survey_data_preprocessing()[0]
    YEARS = [2018, 2019, 2020, 2021]

    sns.set_style("whitegrid")
    plt.rcParams['font.family'] = 'NanumGothic'
    fig, ax = plt.subplots(nrows = 2, ncols = 2, figsize = (15, 10))

    for idx, YEAR in enumerate(YEARS):
        if YEAR == 2018:
            columns = ['세부 카테고리', '자국의 인터넷 사이트/앱', '한국의 여행관련 사이트/앱', '과거여행경험', '친지, 친구, 동료', '관광안내책자']
        else:
            columns = ['세부 카테고리', '한국의 인터넷 사이트/앱', '관광안내책자', '과거여행경험', '친지, 친구, 동료', '관광안내책자']
        for year in [2019, 2020, 2021]:
            if year == 2019:
                info_dict[year] = info_dict[year].rename(columns = {'인터넷 사이트/앱' : '한국의 인터넷 사이트/앱'})
            elif year == 2021:
                info_dict[year] = info_dict[year].rename(columns = {'국내 인터넷/여행관련 사이트/앱' : '한국의 인터넷 사이트/앱'})
            else:
                pass
        data_temp = info_dict[YEAR][columns]
        df_melted = data_temp.melt(id_vars = ['세부 카테고리'], var_name = '수집방법', value_name = '값')
        sns.lineplot(data = df_melted, x = '세부 카테고리', y = '값', hue = '수집방법',
        marker = "o", markersize = 8, ax = ax[idx//2, idx%2])
        ax[idx//2, idx%2].set_title(f'{YEAR}년 정보수집경로')
        ax[idx//2, idx%2].set_xlabel('월')
        ax[idx//2, idx%2].set_ylabel('수집 경로 사용량')
        ax[idx//2, idx%2].legend(title = '카테고리', bbox_to_anchor = (1.05, 1), loc = 'upper left')  # 범례 위치 조정
        ax[idx//2, idx%2].set_xticklabels(ax[idx//2, idx%2].get_xticklabels(), rotation = 45)
        ax[idx//2, idx%2].grid(False)
        fig.tight_layout()
        fig.savefig('./data/streamlit/info_path.png', bbox_inches = 'tight')
        image = Image.open('./data/streamlit/info_path.png')
    
    return image

def survey_data_preprocessing():
    DATA_PATH = './data/물가/설문조사/'
    YEARS = [2018, 2019, 2020, 2021]
    info_dict = {}
    dis_dict = {}
    budget_dict = {}

    for YEAR in YEARS:
        temp_info = pd.read_excel(DATA_PATH + f'{YEAR}년 정보습득경로.xlsx')
        temp_info = survey_drop_row(survey_rename_df(temp_info), YEAR)
        info_dict[YEAR] = temp_info

    for YEAR in YEARS:
        temp_dis = pd.read_excel(DATA_PATH + f'{YEAR}년 여행 시 불만족점.xlsx')
        temp_dis = survey_drop_row(survey_rename_df(temp_dis), YEAR)
        temp_dis['연도'] = YEAR
        dis_dict[YEAR] = temp_dis
    
    for YEAR in YEARS:
        temp_budget = pd.read_excel(DATA_PATH + f'{YEAR}년 제주여행 평가_여행경비.xlsx')
        if YEAR == 2019:
            del temp_budget['Unnamed: 8']
            del temp_budget['Unnamed: 12']
        else:
            pass
        temp_budget.columns = ['분석 카테고리', '세부 카테고리', '사례수', '매우 불만족', '불만족', '보통', '만족', '매우 만족', '매우 불만족 & 불만족', 'mid', '매우 만족 & 만족', '[5점 평균]', '[100점 평균]']
        temp_budget = temp_budget.drop(temp_budget.index[0]).reset_index(drop = True)
        temp_budget = survey_drop_row(temp_budget, YEAR)
        budget_dict[YEAR] = temp_budget
    
    return [info_dict, dis_dict, budget_dict]

def survey_rename_df(df):
    df.iloc[0, 0] = '분석 카테고리'
    df.iloc[0, 1] = '세부 카테고리'
    df.columns = df.iloc[0]
    df = df.drop(df.index[0]).reset_index(drop = True)
    return df

def survey_drop_row(df, YEAR):
    if YEAR == 2020:
        df = df.drop(index = range(0, 3)).reset_index(drop = True)
    else:
        df = df.drop(index = 0).reset_index(drop = True)

    if YEAR == 2021:
        df = df.drop(index = range(9, 44)).reset_index(drop = True)
    elif YEAR == 2020:
        df = df.drop(index = range(7, 42)).reset_index(drop = True)
    else:
        df = df.drop(index = range(12, 45)).reset_index(drop = True)

    return df